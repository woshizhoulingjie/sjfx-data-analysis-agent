import csv
import hashlib
import json
import mimetypes
import os
import re
import stat as stat_module
import time
from datetime import datetime, timezone
from pathlib import Path


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl",
    ".xml", ".html", ".htm", ".log", ".ini", ".cfg", ".yaml", ".yml",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".sql",
}
def _configured_names(variable):
    """Return an explicit, case-insensitive scan exclusion set.

    A forensic inventory must not silently drop hidden, generated, unsupported
    or sensitive entries.  Operators may still exclude known-noise paths, but
    doing so is an explicit policy decision that is surfaced in coverage.
    """
    raw = str(os.getenv(variable, "") or "")
    # The host path separator is the documented delimiter.  Commas and
    # newlines are accepted as a convenience for simple deployment files.
    for delimiter in ("\r", "\n", ","):
        raw = raw.replace(delimiter, os.pathsep)
    return {item.strip().casefold() for item in raw.split(os.pathsep) if item.strip()}


IGNORED_DIRS = _configured_names("SCAN_IGNORED_DIRS")
IGNORED_FILES = _configured_names("SCAN_IGNORED_FILES")
SENSITIVE_EXTENSIONS = {".key", ".pem", ".p12", ".pfx", ".keystore"}
SENSITIVE_FILENAMES = {
    ".env", ".netrc", ".npmrc", ".pypirc", "credentials", "credentials.json",
    "id_rsa", "id_dsa", "id_ecdsa", "id_ed25519",
}


def should_ignore_file(name):
    return str(name or "").casefold() in IGNORED_FILES


def is_sensitive_file(name):
    """Identify files whose presence is inventoried but content is restricted."""
    lower_name = str(name or "").casefold()
    return (
        lower_name in SENSITIVE_FILENAMES
        or lower_name.startswith(".env.")
        or Path(lower_name).suffix in SENSITIVE_EXTENSIONS
    )


def natural_key(value):
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", value)]


def now_iso():
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def path_id(path):
    return hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:16]


def human_size(value):
    size = float(value)
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if size < 1024 or unit == "TB":
            return "{:.1f} {}".format(size, unit)
        size /= 1024


def resolve_under(root, requested):
    root_path = Path(root).resolve()
    candidate = Path(requested)
    if not candidate.is_absolute():
        candidate = root_path / candidate
    candidate = candidate.resolve()
    try:
        candidate.relative_to(root_path)
    except ValueError:
        raise ValueError("请求路径超出当前扫描根目录")
    return candidate


def _file_metadata(path, root, stat_result=None):
    # The caller normally supplies DirEntry.lstat() data. Reopening a path
    # here would create a type-check/open race if a writable NAS entry is
    # swapped for a symlink between directory enumeration and metadata read.
    stat = stat_result if stat_result is not None else path.stat()
    mime, _ = mimetypes.guess_type(str(path))
    rel = str(path.relative_to(root)).replace("\\", "/")
    metadata = {
        "id": path_id(path),
        "name": path.name,
        "path": rel,
        "kind": "file",
        "extension": path.suffix.lower(),
        "mime_type": mime or "application/octet-stream",
        "size": stat.st_size,
        "size_human": human_size(stat.st_size),
        "modified_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).replace(microsecond=0).isoformat(),
        # Nanosecond precision is part of the checkpoint precondition.  The
        # ISO timestamp is intentionally human-friendly and loses precision.
        "modified_at_ns": int(stat.st_mtime_ns),
        "device": int(stat.st_dev),
        "inode": int(stat.st_ino),
    }
    if is_sensitive_file(path.name):
        metadata.update({
            "sensitive": True,
            "content_analysis_allowed": False,
            "content_policy": "metadata_only_sensitive",
        })
    return metadata


def _symlink_metadata(entry, root, item_path=None):
    """Inventory a link itself without resolving or following its target."""
    # ``os.scandir(directory_fd)`` returns entries whose ``path`` is only the
    # basename on POSIX. Use the caller's absolute lexical path so inventory
    # IDs and relative paths remain anchored to the scan root.
    path = Path(item_path) if item_path is not None else Path(entry.path)
    stat = entry.stat(follow_symlinks=False)
    rel = str(path.relative_to(root)).replace("\\", "/")
    return {
        "id": path_id(path),
        "name": entry.name,
        "path": rel,
        "kind": "symlink",
        "extension": path.suffix.lower(),
        "size": 0,
        "link_metadata_size": int(stat.st_size),
        "size_human": human_size(0),
        "modified_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).replace(microsecond=0).isoformat(),
        "modified_at_ns": int(stat.st_mtime_ns),
        "content_analysis_allowed": False,
        "content_policy": "inventory_only_symlink_target_not_followed",
    }


def scan_directory(root, max_files=10000, max_depth=32, progress_callback=None,
                   activity_callback=None, cancel_check=None,
                   max_directories=None, max_nodes=None,
                   max_entries_per_directory=None):
    """Build a bounded physical inventory without following symbolic links.

    File, directory, total-node, depth and per-directory entry bounds protect
    the web/worker process from pathological media. Reaching any bound is
    visible in the result instead of being silently treated as complete. The
    optional callbacks let the worker publish progress and honour cancellation
    without coupling this module to the task database.
    """
    root = Path(root).expanduser().resolve()
    if not root.exists() or not root.is_dir():
        raise ValueError("目录不存在或不是文件夹")
    max_files = max(1, min(int(max_files), 1_000_000))
    max_depth = max(1, min(int(max_depth), 256))
    max_directories = max(
        1,
        min(1_000_000, int(
            os.getenv("MAX_SCAN_DIRECTORIES", "50000")
            if max_directories is None else max_directories
        )),
    )
    max_nodes = max(
        2,
        min(2_000_000, int(
            os.getenv("MAX_SCAN_NODES", "100001")
            if max_nodes is None else max_nodes
        )),
    )
    max_entries_per_directory = max(
        1,
        min(250000, int(
            os.getenv("MAX_SCAN_ENTRIES_PER_DIRECTORY", "50000")
            if max_entries_per_directory is None else max_entries_per_directory
        )),
    )
    count = 0
    total_size = 0
    type_counts = {}
    errors = []
    truncated = False
    ignored_file_count = 0
    ignored_directory_count = 0
    skipped_symlink_count = 0
    depth_limited_directory_count = 0
    visited_directory_count = 0
    scanned_node_count = 0
    directory_limited_count = 0
    node_limited_count = 0
    entry_limited_directory_count = 0
    last_activity_at = 0.0

    def publish_activity(current_path, force=False):
        nonlocal last_activity_at
        if not activity_callback:
            return
        now = time.monotonic()
        if force or now - last_activity_at >= 0.75:
            activity_callback(count, visited_directory_count, str(current_path))
            last_activity_at = now

    def walk(folder, depth=0, folder_fd=None, folder_stat=None):
        nonlocal count, total_size, truncated, ignored_file_count, ignored_directory_count
        nonlocal skipped_symlink_count, depth_limited_directory_count, visited_directory_count
        nonlocal scanned_node_count, directory_limited_count, node_limited_count
        nonlocal entry_limited_directory_count
        visited_directory_count += 1
        scanned_node_count += 1
        publish_activity(folder, force=visited_directory_count == 1)
        node = {
            "id": path_id(folder),
            "name": folder.name or str(folder),
            "path": str(folder.relative_to(root)).replace("\\", "/") if folder != root else ".",
            "kind": "directory",
            "children": [],
            "file_count": 0,
            "direct_file_count": 0,
            "directory_count": 0,
            "direct_directory_count": 0,
            "total_size": 0,
            "type_counts": {},
            "scan_depth": depth,
        }
        if folder_stat is not None:
            node["device"] = int(folder_stat.st_dev)
            node["inode"] = int(folder_stat.st_ino)
        if cancel_check:
            cancel_check()
        if depth >= max_depth:
            depth_limited_directory_count += 1
            node["depth_limited"] = True
            node["simple_summary"] = "目录层级超过安全上限，未继续向下扫描。"
            return node
        try:
            entries = []
            # On POSIX, scan an already-open directory object. Child
            # directories are opened relative to this descriptor with
            # O_NOFOLLOW, so renaming an entry to a symlink cannot redirect
            # traversal outside the imported root.
            scan_target = folder_fd if folder_fd is not None else str(folder)
            with os.scandir(scan_target) as iterator:
                for entry_index, entry in enumerate(iterator, 1):
                    if cancel_check:
                        cancel_check()
                    if entry_index > max_entries_per_directory:
                        truncated = True
                        entry_limited_directory_count += 1
                        node["entry_limited"] = True
                        break
                    entries.append(entry)
                    if entry_index % 64 == 0:
                        publish_activity(folder)
            if cancel_check:
                cancel_check()
            entries.sort(
                key=lambda e: (not e.is_dir(follow_symlinks=False), natural_key(e.name))
            )
        except (OSError, PermissionError) as exc:
            errors.append({"path": str(folder), "error": str(exc)})
            return node
        for entry in entries:
            if cancel_check:
                cancel_check()
            if entry.name.casefold() in IGNORED_DIRS:
                ignored_directory_count += 1
                continue
            if count >= max_files:
                truncated = True
                break
            if scanned_node_count >= max_nodes:
                truncated = True
                node_limited_count += 1
                node["node_limited"] = True
                break
            try:
                item_path = folder / entry.name
                if entry.is_symlink():
                    skipped_symlink_count += 1
                    node["children"].append(_symlink_metadata(entry, root, item_path))
                    scanned_node_count += 1
                    continue
                entry_stat = entry.stat(follow_symlinks=False)
                if stat_module.S_ISDIR(entry_stat.st_mode):
                    if visited_directory_count >= max_directories:
                        truncated = True
                        directory_limited_count += 1
                        node["directory_limited"] = True
                        continue
                    child_fd = None
                    try:
                        if folder_fd is not None:
                            directory_flags = (
                                os.O_RDONLY
                                | getattr(os, "O_DIRECTORY", 0)
                                | getattr(os, "O_NOFOLLOW", 0)
                                | getattr(os, "O_CLOEXEC", 0)
                            )
                            child_fd = os.open(entry.name, directory_flags, dir_fd=folder_fd)
                            opened_stat = os.fstat(child_fd)
                            if (
                                int(opened_stat.st_dev) != int(entry_stat.st_dev)
                                or int(opened_stat.st_ino) != int(entry_stat.st_ino)
                            ):
                                raise OSError("目录项在扫描期间发生替换")
                            child = walk(item_path, depth + 1, child_fd, opened_stat)
                        else:
                            # Windows fallback: reject a reparse/symlink at the
                            # last possible point; parse-time handle validation
                            # provides the authoritative containment check.
                            if item_path.is_symlink():
                                raise OSError("目录项在扫描期间变为符号链接")
                            child = walk(item_path, depth + 1, None, entry_stat)
                    finally:
                        if child_fd is not None:
                            os.close(child_fd)
                    node["children"].append(child)
                    node["direct_directory_count"] += 1
                    node["directory_count"] += 1 + child["directory_count"]
                    node["file_count"] += child["file_count"]
                    node["total_size"] += child["total_size"]
                    for ext, value in child["type_counts"].items():
                        node["type_counts"][ext] = node["type_counts"].get(ext, 0) + value
                elif stat_module.S_ISREG(entry_stat.st_mode):
                    if should_ignore_file(entry.name):
                        ignored_file_count += 1
                        continue
                    meta = _file_metadata(item_path, root, entry_stat)
                    node["children"].append(meta)
                    scanned_node_count += 1
                    count += 1
                    total_size += meta["size"]
                    node["file_count"] += 1
                    node["direct_file_count"] += 1
                    node["total_size"] += meta["size"]
                    ext = meta["extension"] or "[无扩展名]"
                    type_counts[ext] = type_counts.get(ext, 0) + 1
                    node["type_counts"][ext] = node["type_counts"].get(ext, 0) + 1
                    if progress_callback and (count == 1 or count % 25 == 0):
                        progress_callback(count)
                    publish_activity(item_path)
            except (OSError, PermissionError) as exc:
                errors.append({"path": str(folder / entry.name), "error": str(exc)})
        node["size_human"] = human_size(node["total_size"])
        node["type_counts"] = dict(sorted(node["type_counts"].items(), key=lambda item: (-item[1], item[0])))
        top_types = list(node["type_counts"].items())[:4]
        if node["file_count"]:
            type_text = "、".join("{} {}个".format(ext, value) for ext, value in top_types)
            node["simple_summary"] = (
                "本文件夹当前层有 {direct_files} 个文件、{direct_dirs} 个子目录；"
                "递归范围共 {files} 个文件、{dirs} 个目录，总大小 {size}。主要类型：{types}。"
            ).format(
                direct_files=node["direct_file_count"], direct_dirs=node["direct_directory_count"],
                files=node["file_count"], dirs=node["directory_count"], size=node["size_human"],
                types=type_text or "无扩展名统计",
            )
        else:
            node["simple_summary"] = "本文件夹当前为空，未发现可分析文件。"
        publish_activity(folder)
        return node

    root_fd = None
    root_stat = None
    try:
        if os.name != "nt" and os.open in getattr(os, "supports_dir_fd", set()):
            root_fd = os.open(
                str(root),
                os.O_RDONLY
                | getattr(os, "O_DIRECTORY", 0)
                | getattr(os, "O_NOFOLLOW", 0)
                | getattr(os, "O_CLOEXEC", 0),
            )
            root_stat = os.fstat(root_fd)
        elif root.exists():
            root_stat = root.stat()
        tree = walk(root, folder_fd=root_fd, folder_stat=root_stat)
    finally:
        if root_fd is not None:
            os.close(root_fd)
    publish_activity(root, force=True)
    return {
        "root": str(root),
        "scanned_at": now_iso(),
        "file_count": count,
        "directory_count": tree["directory_count"],
        "ignored_file_count": ignored_file_count,
        "ignored_directory_count": ignored_directory_count,
        "ignore_policy": {
            "directories": sorted(IGNORED_DIRS),
            "files": sorted(IGNORED_FILES),
            "default_is_full_inventory": not IGNORED_DIRS and not IGNORED_FILES,
        },
        "symlink_count": skipped_symlink_count,
        "skipped_symlink_count": skipped_symlink_count,
        "depth_limited_directory_count": depth_limited_directory_count,
        "directory_limited_count": directory_limited_count,
        "node_limited_count": node_limited_count,
        "entry_limited_directory_count": entry_limited_directory_count,
        "scanned_directory_count": visited_directory_count,
        "scanned_node_count": scanned_node_count,
        "max_depth": max_depth,
        "max_files": max_files,
        "max_directories": max_directories,
        "max_nodes": max_nodes,
        "max_entries_per_directory": max_entries_per_directory,
        "total_size": total_size,
        "total_size_human": human_size(total_size),
        "type_counts": dict(sorted(type_counts.items(), key=lambda item: (-item[1], item[0]))),
        "truncated": truncated,
        "scan_error_count": len(errors),
        "errors": errors[:100],
        "tree": tree,
    }


def scan_inventory_slice(root, cursor=None, slice_entries=1000, slice_seconds=20,
                         max_depth=32, max_files=1_000_000, max_directories=1_000_000,
                         max_nodes=2_000_000, cancel_check=None,
                         yield_check=None, activity_callback=None,
                         manifest_dir=None):
    """Enumerate one durable, lexically ordered inventory slice.

    When ``manifest_dir`` is supplied, each directory is enumerated and sorted
    exactly once into a durable local manifest. Subsequent slices continue from
    a byte offset instead of rescanning and resorting a huge flat directory.
    """
    root = Path(root).expanduser().resolve()
    if not root.exists() or not root.is_dir():
        raise ValueError("目录不存在或不是文件夹")
    max_depth = max(1, min(256, int(max_depth or 32)))
    max_files = max(1, min(1_000_000, int(max_files or 1_000_000)))
    max_directories = max(1, min(1_000_000, int(max_directories or 1_000_000)))
    max_nodes = max(2, min(2_000_000, int(max_nodes or 2_000_000)))
    slice_entries = max(1, min(10000, int(slice_entries or 1000)))
    slice_seconds = max(1.0, min(300.0, float(slice_seconds or 20)))
    started = time.monotonic()
    manifest_root = Path(manifest_dir).expanduser().resolve() if manifest_dir else None
    if manifest_root is not None:
        manifest_root.mkdir(parents=True, exist_ok=True)

    state = dict(cursor or {})
    if state and state.get("version") != 1:
        raise ValueError("扫描游标版本不兼容")
    if state and str(state.get("root") or "") != str(root):
        raise ValueError("扫描游标与当前目录不匹配")
    records = []
    if not state:
        root_stat = root.stat()
        root_node = {
            "id": path_id(root), "name": root.name or str(root), "path": ".",
            "kind": "directory", "scan_depth": 0,
            "device": int(root_stat.st_dev), "inode": int(root_stat.st_ino),
        }
        records.append({"path": ".", "parent_path": None, "position": 0, "payload": root_node})
        state = {
            "version": 1, "root": str(root),
            "stack": [{"path": ".", "depth": 0, "after_key": None}],
            "file_count": 0, "directory_count": 1, "node_count": 1,
            "symlink_count": 0, "ignored_file_count": 0,
            "ignored_directory_count": 0, "depth_limited_directory_count": 0,
            "total_size": 0, "type_counts": {}, "errors": [],
        }

    def relative_path(parent_path, name):
        return name if parent_path == "." else parent_path.rstrip("/") + "/" + name

    def should_yield():
        return (
            len(records) >= slice_entries
            or time.monotonic() - started >= slice_seconds
            or (yield_check is not None and bool(yield_check()))
        )

    def next_manifest_item(frame, directory_path):
        manifest_path = frame.get("manifest_path")
        if not manifest_path:
            try:
                with os.scandir(str(directory_path)) as iterator:
                    names = [entry.name for entry in iterator]
            except (OSError, PermissionError) as exc:
                state["errors"].append({"path": frame["path"], "error": str(exc)[:1000]})
                return None
            names.sort(key=lambda value: (value.casefold(), value))
            after = tuple(frame.get("after_key") or ())
            digest = hashlib.sha256(
                (str(root) + "\0" + str(frame["path"])).encode("utf-8", errors="replace")
            ).hexdigest()
            target = manifest_root / (digest + ".jsonl")
            temporary = manifest_root / (digest + ".{}.tmp".format(os.getpid()))
            with temporary.open("w", encoding="utf-8", newline="\n") as handle:
                for position, name in enumerate(names):
                    key = (name.casefold(), name)
                    if after and key <= after:
                        continue
                    handle.write(json.dumps([position, name], ensure_ascii=False) + "\n")
                handle.flush()
                os.fsync(handle.fileno())
            os.replace(str(temporary), str(target))
            frame["manifest_path"] = str(target)
            frame["manifest_offset"] = 0
            manifest_path = str(target)
        manifest_path = Path(manifest_path).resolve()
        try:
            manifest_path.relative_to(manifest_root)
        except ValueError as exc:
            raise ValueError("扫描清单游标指向了状态目录之外") from exc
        try:
            with manifest_path.open("rb") as handle:
                handle.seek(max(0, int(frame.get("manifest_offset") or 0)))
                line = handle.readline()
                frame["manifest_offset"] = handle.tell()
        except FileNotFoundError:
            # A state-disk cleanup invalidates only the directory manifest. It
            # is safe to rebuild from after_key without losing saved records.
            frame.pop("manifest_path", None)
            frame.pop("manifest_offset", None)
            return next_manifest_item(frame, directory_path)
        if not line:
            try:
                manifest_path.unlink()
            except OSError:
                pass
            return None
        position, name = json.loads(line.decode("utf-8"))
        key = (str(name).casefold(), str(name))
        return int(position), str(name), key

    while state["stack"]:
        if cancel_check is not None:
            cancel_check()
        frame = state["stack"][-1]
        directory_path = root if frame["path"] == "." else root / frame["path"]
        if manifest_root is not None:
            next_item = next_manifest_item(frame, directory_path)
        else:
            try:
                with os.scandir(str(directory_path)) as iterator:
                    names = [entry.name for entry in iterator]
            except (OSError, PermissionError) as exc:
                state["errors"].append({"path": frame["path"], "error": str(exc)[:1000]})
                state["stack"].pop()
                continue
            names.sort(key=lambda value: (value.casefold(), value))
            after = tuple(frame.get("after_key") or ())
            next_item = None
            for position, name in enumerate(names):
                key = (name.casefold(), name)
                if not after or key > after:
                    next_item = (position, name, key)
                    break
        if next_item is None:
            state["stack"].pop()
            continue

        position, name, key = next_item
        frame["after_key"] = list(key)
        child_path = relative_path(frame["path"], name)
        item_path = directory_path / name
        try:
            item_stat = item_path.lstat()
            mode = item_stat.st_mode
            if stat_module.S_ISLNK(mode):
                payload = {
                    "id": path_id(item_path), "name": name, "path": child_path,
                    "kind": "symlink", "extension": item_path.suffix.lower(),
                    "size": 0, "link_metadata_size": int(item_stat.st_size),
                    "size_human": human_size(0),
                    "modified_at": datetime.fromtimestamp(
                        item_stat.st_mtime, timezone.utc
                    ).replace(microsecond=0).isoformat(),
                    "modified_at_ns": int(item_stat.st_mtime_ns),
                    "content_analysis_allowed": False,
                    "content_policy": "inventory_only_symlink_target_not_followed",
                }
                records.append({
                    "path": child_path, "parent_path": frame["path"],
                    "position": position, "payload": payload,
                })
                state["symlink_count"] += 1
                state["node_count"] += 1
            elif stat_module.S_ISDIR(mode):
                if name.casefold() in IGNORED_DIRS:
                    state["ignored_directory_count"] += 1
                    continue
                if state["directory_count"] >= max_directories or state["node_count"] >= max_nodes:
                    raise ValueError("目录清单超过已配置的安全上限，扫描未完成")
                depth = int(frame["depth"]) + 1
                payload = {
                    "id": path_id(item_path), "name": name, "path": child_path,
                    "kind": "directory", "scan_depth": depth,
                    "device": int(item_stat.st_dev), "inode": int(item_stat.st_ino),
                }
                if depth >= max_depth:
                    payload["depth_limited"] = True
                    payload["simple_summary"] = "目录层级超过安全上限，未继续向下扫描。"
                    state["depth_limited_directory_count"] += 1
                records.append({
                    "path": child_path, "parent_path": frame["path"],
                    "position": position, "payload": payload,
                })
                state["directory_count"] += 1
                state["node_count"] += 1
                if depth < max_depth:
                    state["stack"].append({
                        "path": child_path, "depth": depth, "after_key": None,
                    })
            elif stat_module.S_ISREG(mode):
                if should_ignore_file(name):
                    state["ignored_file_count"] += 1
                    continue
                if state["file_count"] >= max_files or state["node_count"] >= max_nodes:
                    raise ValueError("文件清单超过已配置的安全上限，扫描未完成")
                payload = _file_metadata(item_path, root, item_stat)
                records.append({
                    "path": child_path, "parent_path": frame["path"],
                    "position": position, "payload": payload,
                })
                state["file_count"] += 1
                state["node_count"] += 1
                state["total_size"] += int(payload.get("size") or 0)
                extension = payload.get("extension") or "[无扩展名]"
                state["type_counts"][extension] = state["type_counts"].get(extension, 0) + 1
        except ValueError:
            raise
        except (OSError, PermissionError) as exc:
            state["errors"].append({"path": child_path, "error": str(exc)[:1000]})

        if activity_callback is not None:
            activity_callback(
                state["file_count"], state["directory_count"], child_path
            )
        if should_yield():
            break

    complete = not state["stack"]
    state["type_counts"] = dict(sorted(
        state["type_counts"].items(), key=lambda item: (-item[1], item[0])
    ))
    return {"records": records, "cursor": state, "complete": complete}


def _decode_bytes(data):
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            pass
    return data.decode("utf-8", errors="replace")


def extract_text(path, max_chars=60000):
    path = Path(path)
    ext = path.suffix.lower()
    text = ""
    parser = "unsupported"
    warnings = []
    metadata = {}
    try:
        if ext in TEXT_EXTENSIONS:
            with path.open("rb") as stream:
                raw = stream.read(max_chars * 4 + 1)
            text = _decode_bytes(raw[: max_chars * 4])
            if len(raw) > max_chars * 4:
                warnings.append("文本文件超过本地完整读取上限")
            parser = "text"
        elif ext == ".pdf":
            try:
                try:
                    from pypdf import PdfReader
                except ImportError:
                    from PyPDF2 import PdfReader
                reader = PdfReader(str(path))
                metadata["page_count"] = len(reader.pages)
                parts = []
                extracted_pages = 0
                empty_pages = 0
                accumulated = 0
                for page_number, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        extracted_pages += 1
                    else:
                        empty_pages += 1
                    section = "[第 {} 页]\n{}".format(page_number, page_text)
                    parts.append(section)
                    accumulated += len(section)
                    if accumulated >= max_chars:
                        break
                text = "\n\n".join(parts)
                metadata["processed_pages"] = len(parts)
                metadata["text_pages"] = extracted_pages
                metadata["empty_text_pages"] = empty_pages
                metadata["pages_omitted_by_limit"] = max(0, len(reader.pages) - len(parts))
                parser = "PyPDF2"
            except ImportError:
                warnings.append("未安装 PyPDF2，当前只能读取文件元数据")
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                text = "\n".join(p.text for p in doc.paragraphs)
                metadata["paragraph_count"] = len(doc.paragraphs)
                metadata["table_count"] = len(doc.tables)
                parser = "python-docx"
            except ImportError:
                warnings.append("未安装 python-docx，当前只能读取文件元数据")
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                presentation = Presentation(str(path))
                parts = []
                for index, slide in enumerate(presentation.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    parts.append("[幻灯片 {}]\n{}".format(index, "\n".join(slide_text)))
                text = "\n\n".join(parts)
                metadata["slide_count"] = len(presentation.slides)
                parser = "python-pptx"
            except ImportError:
                warnings.append("未安装 python-pptx，当前只能读取文件元数据")
        elif ext in (".xlsx", ".xlsm"):
            try:
                from openpyxl import load_workbook
                workbook = load_workbook(str(path), read_only=True, data_only=True)
                parts = []
                accumulated = 0
                stopped = False
                for sheet in workbook.worksheets:
                    parts.append("[工作表 {}]".format(sheet.title))
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join("" if value is None else str(value) for value in row)
                        parts.append(row_text)
                        accumulated += len(row_text)
                        if accumulated >= max_chars:
                            stopped = True
                            break
                    if stopped:
                        break
                text = "\n".join(parts)
                metadata["sheet_count"] = len(workbook.worksheets)
                parser = "openpyxl"
            except ImportError:
                warnings.append("未安装 openpyxl，当前只能读取文件元数据")
        elif ext in {".doc", ".xls", ".ppt"}:
            metadata["capability"] = "out_of_scope_legacy_office"
            warnings.append("旧版 Office 格式当前仅登记元数据；正文解析需要 LibreOffice/专用转换器")
        else:
            warnings.append("该文件类型尚未配置正文解析器")
    except Exception as exc:
        warnings.append("解析失败：{}".format(exc))
    original_length = len(text)
    if original_length > max_chars:
        text = text[:max_chars]
        warnings.append("正文已截断：原始提取字符数 {}".format(original_length))
    return {
        "text": text.strip(),
        "parser": parser,
        "warnings": warnings,
        "metadata": metadata,
        "char_count": original_length,
        "truncated": original_length > max_chars,
    }


def folder_context(folder, root, max_files=30, max_chars=50000, max_depth=32):
    folder = Path(folder)
    root = Path(root)
    inventory = []
    excerpts = []
    sampled = 0
    total_files = 0
    total_dirs = 0
    total_size = 0
    type_counts = {}
    candidates = []
    root_depth = len(folder.resolve().parts)
    for current_root, dirs, files in os.walk(str(folder), followlinks=False):
        current = Path(current_root)
        depth = max(0, len(current.resolve().parts) - root_depth)
        dirs[:] = [
            d for d in dirs
            if d.casefold() not in IGNORED_DIRS and not (current / d).is_symlink()
        ]
        if depth >= max(1, int(max_depth)):
            dirs[:] = []
        total_dirs += len(dirs)
        for name in sorted(files, key=natural_key):
            if should_ignore_file(name):
                continue
            path = Path(current_root) / name
            try:
                meta = _file_metadata(path, root)
                total_files += 1
                total_size += meta["size"]
                ext = meta["extension"] or "[无扩展名]"
                type_counts[ext] = type_counts.get(ext, 0) + 1
                candidates.append((path, meta))
            except (OSError, PermissionError):
                continue
    if len(candidates) <= max_files:
        selected = candidates
    elif max_files <= 1:
        selected = candidates[:1]
    else:
        indices = sorted(set(int(round(index * (len(candidates) - 1) / float(max_files - 1))) for index in range(max_files)))
        selected = [candidates[index] for index in indices]
    documents = []
    remaining_chars = max_chars
    per_file_budget = min(6000, max(800, max_chars // max(1, len(selected))))
    for path, meta in selected:
        if remaining_chars <= 0:
            break
        per_file_limit = min(per_file_budget, remaining_chars)
        extracted = extract_text(path, max_chars=per_file_limit)
        text_sample = extracted["text"][:per_file_limit]
        inventory.append({"path": meta["path"], "size": meta["size"], "extension": meta["extension"]})
        documents.append({
            "path": meta["path"], "extension": meta["extension"], "text": text_sample,
            "parser": extracted["parser"], "warnings": extracted["warnings"],
        })
        if text_sample:
            excerpt = "### {}\n{}".format(meta["path"], text_sample)
            excerpts.append(excerpt)
            remaining_chars -= len(excerpt)
        sampled += 1
    joined = "\n\n".join(excerpts)
    return {
        "inventory": inventory,
        "excerpts": joined[:max_chars],
        "sampled_files": sampled,
        "total_files": total_files,
        "total_dirs": total_dirs,
        "total_size": total_size,
        "total_size_human": human_size(total_size),
        "type_counts": dict(sorted(type_counts.items(), key=lambda item: (-item[1], item[0]))),
        "sample_truncated": total_files > sampled,
        "documents": documents,
    }
